"""Gera Apresentacao-UrbanData-Rota-Eletronica.pptx (5 slides)."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# UrbanData palette
BG_DARK = RGBColor(0x00, 0x1B, 0x28)
GREEN = RGBColor(0x20, 0xB5, 0x73)
PETROL = RGBColor(0x1C, 0x72, 0x7A)
TEXT_LIGHT = RGBColor(0xD0, 0xD0, 0xD0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
SUBTLE = RGBColor(0xA8, 0xB8, 0xC0)

OUTPUT = Path(r"C:\Users\Jessi\Downloads\Apresentacao-UrbanData-Rota-Eletronica.pptx")


def set_slide_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_bar(slide, prs: Presentation) -> None:
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0),
        Inches(0),
        prs.slide_width,
        Inches(0.12),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()


def style_title(tf, size_pt: int = 32, color: RGBColor = WHITE) -> None:
    p = tf.paragraphs[0]
    p.font.size = Pt(size_pt)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT


def add_bullets(text_frame, items: list[str], size_pt: int = 18) -> None:
    text_frame.clear()
    for i, item in enumerate(items):
        p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size_pt)
        p.font.color.rgb = TEXT_LIGHT
        p.space_after = Pt(10)
        p.bullet = True


def add_table_slide(slide, prs, headers: list[str], rows: list[list[str]]) -> None:
    cols, row_count = len(headers), len(rows) + 1
    left, top = Inches(0.6), Inches(2.0)
    width, height = Inches(8.8), Inches(0.45 * row_count)
    table = slide.shapes.add_table(row_count, cols, left, top, width, height).table

    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.color.rgb = WHITE

    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = TEXT_LIGHT

    for r in range(row_count):
        for c in range(cols):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PETROL if r == 0 else RGBColor(0x00, 0x25, 0x35)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # Slide 1 — Título
    s1 = prs.slides.add_slide(blank)
    set_slide_bg(s1, BG_DARK)
    add_title_bar(s1, prs)
    box = s1.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].text = "UrbanData — Rota Eletrônica Escolar"
    style_title(tf, 36, GREEN)
    p2 = tf.add_paragraph()
    p2.text = "Gestão integrada do transporte escolar municipal"
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_LIGHT
    p2.space_before = Pt(16)
    p3 = tf.add_paragraph()
    p3.text = "UrbanData · Painel web + API · PostgreSQL"
    p3.font.size = Pt(14)
    p3.font.color.rgb = SUBTLE
    p3.space_before = Pt(24)

    # Slide 2 — Problema
    s2 = prs.slides.add_slide(blank)
    set_slide_bg(s2, BG_DARK)
    add_title_bar(s2, prs)
    t2 = s2.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(0.8))
    t2.text_frame.text = "O desafio"
    style_title(t2.text_frame, 28, GREEN)
    b2 = s2.shapes.add_textbox(Inches(0.6), Inches(1.5), Inches(8.8), Inches(4.5))
    add_bullets(
        b2.text_frame,
        [
            "Frota, alunos e rotas em planilhas e processos manuais",
            "Dificuldade para planejar trajetos e turnos (manhã / tarde)",
            "Pouca visibilidade do que está rodando hoje (atrasos, veículos, paradas)",
        ],
        20,
    )

    # Slide 3 — Solução (tabela)
    s3 = prs.slides.add_slide(blank)
    set_slide_bg(s3, BG_DARK)
    add_title_bar(s3, prs)
    t3 = s3.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(0.8))
    t3.text_frame.text = "O que o sistema faz"
    style_title(t3.text_frame, 28, GREEN)
    add_table_slide(
        s3,
        prs,
        ["Área", "Função"],
        [
            ["Cadastros", "Municípios, escolas, garagens, veículos, motoristas, alunos"],
            ["Roteirização", "Montagem de rotas com otimização Google Maps"],
            ["Escalas", "Programação rota + veículo + motorista por data"],
            ["Monitoramento", "Mapa e dashboard da operação do dia"],
        ],
    )

    # Slide 4 — Como funciona
    s4 = prs.slides.add_slide(blank)
    set_slide_bg(s4, BG_DARK)
    add_title_bar(s4, prs)
    t4 = s4.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(0.8))
    t4.text_frame.text = "Como funciona"
    style_title(t4.text_frame, 28, GREEN)
    b4 = s4.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(8.8), Inches(3.2))
    add_bullets(
        b4.text_frame,
        [
            "Cadastrar — município, escolas, frota e alunos (importação em lote)",
            "Roteirizar — garagem → paradas → escola e otimizar trajeto",
            "Escalar — veículo e motorista para cada dia",
            "Acompanhar — dashboard e mapa de monitoramento",
        ],
        19,
    )
    note = s4.shapes.add_textbox(Inches(0.6), Inches(5.0), Inches(8.8), Inches(1.2))
    note.text_frame.word_wrap = True
    p = note.text_frame.paragraphs[0]
    p.text = (
        "Controle de acesso: Administrador, Gestor e Operador — "
        "cada perfil vê apenas os municípios permitidos."
    )
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = PETROL

    # Slide 5 — Entrega e visão
    s5 = prs.slides.add_slide(blank)
    set_slide_bg(s5, BG_DARK)
    add_title_bar(s5, prs)
    t5 = s5.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(0.8))
    t5.text_frame.text = "Entrega atual e visão"
    style_title(t5.text_frame, 28, GREEN)

    hoje = s5.shapes.add_textbox(Inches(0.6), Inches(1.35), Inches(4.2), Inches(0.5))
    hoje.text_frame.text = "Hoje"
    style_title(hoje.text_frame, 20, GREEN)

    b_hoje = s5.shapes.add_textbox(Inches(0.6), Inches(1.85), Inches(4.3), Inches(2.8))
    add_bullets(
        b_hoje.text_frame,
        [
            "Painel web operacional",
            "API com dados no PostgreSQL (UrbanData)",
            "Integração Google Maps e ViaCEP",
        ],
        17,
    )

    roadmap = s5.shapes.add_textbox(Inches(5.2), Inches(1.35), Inches(4.2), Inches(0.5))
    roadmap.text_frame.text = "Próximo passo"
    style_title(roadmap.text_frame, 20, GREEN)

    b_road = s5.shapes.add_textbox(Inches(5.2), Inches(1.85), Inches(4.3), Inches(2.8))
    add_bullets(
        b_road.text_frame,
        [
            "App Motorista — rota ativa, paradas, intercorrências",
            "App Pais — itinerário e status do transporte",
        ],
        17,
    )

    footer = s5.shapes.add_textbox(Inches(0.6), Inches(5.8), Inches(8.8), Inches(0.6))
    footer.text_frame.paragraphs[0].text = "www.urbandata.com.br"
    footer.text_frame.paragraphs[0].font.size = Pt(12)
    footer.text_frame.paragraphs[0].font.color.rgb = SUBTLE
    footer.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Arquivo criado: {path}")
